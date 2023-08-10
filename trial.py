from matplotlib import figure
from pptx import Presentation
import pandas as pd
import matplotlib.pyplot as plt
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

from powerpoint import analysis_week, plot_split



# Define days dictionary
dict_days = {
    'MON': 'Monday',
    'TUE': 'Tuesday',
    'WED': 'Wednesday',
    'THU': 'Thursday',
    'FRI': 'Friday',
    'SAT': 'Saturday',
    'SUN': 'Sunday'
}
# Weeks
LIST_WEEKS = dict_days['WEEK'].unique()

# Load dataframes
df_day = pd.read_csv('volumes_per_day.csv', index_col=0)
df_lior = pd.read_csv('lines_per_day.csv', index_col=0)
COLS_IN = list(df_lior.columns[0:8])

# Create PowerPoint presentation
prs = Presentation()

# Introduction Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 32, 96)
title.text = "WAREHOUSE WORKLOAD ANALYSIS"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
subtitle.text = f'Orders/day for the last {len(LIST_WEEKS)} weeks'
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Analysis per Week
image_slide_layout = prs.slide_layouts[5]
page = 1

for WEEK in LIST_WEEKS:
    avg_ratio, max_ratio, busy_day, max_lines, total_lines = analysis_week(df_day, WEEK)
    
    slide = prs.slides.add_slide(image_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = f'Warehouse Workload ({WEEK})'
    
    left = Inches(0.75)
    top = Inches(1.25)
    pic = slide.shapes.add_picture(f'{WEEK}.png', left, top, height=Inches(4.5))
    
    left = Inches(0.75)
    top = Inches(1.5) + Inches(4)
    width = Inches(9)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = 'Analysis'
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = f'• {total_lines} have been prepared during the week'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f'• {busy_day} has been the busiest day with {max_lines} prepared'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f'• {avg_ratio} on average with a maximum of {max_ratio}'
    p.level = 1
    
    filename = f'{WEEK}.png'
    path_plot = filename
    figure.savefig(path_plot, dpi=figure.Figure.dpi)
    
    page += 1

# Order Profile Slide
total_orders, LIST_ANALYSIS = plot_split(df_lior)

slide = prs.slides.add_slide(image_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'Order Profile'

left = Inches(0.75)
top = Inches(1.25)
pic = slide.shapes.add_picture('SPLIT.png', left, top, height=Inches(4.5))

left = Inches(0.75)
top = Inches(1.5) + Inches(4)
width = Inches(9)
height = Inches(2)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = f'{total_orders} prepared'
p.font.size = Pt(18)

for l in LIST_ANALYSIS:
    p = tf.add_paragraph()
    p.text = f'• {l}'
    p.level = 1

txBox = slide.shapes.add_textbox(Inches(9), Inches(6.75), Inches(1), Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = str('{}/{}'.format(page, len(LIST_WEEKS) + 1))
p.font.size = Pt(15)
page += 1

# Save the presentation
prs.save('Warehouse_Workload_Report.pptx')
