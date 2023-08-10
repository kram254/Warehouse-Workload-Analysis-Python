from pptx import Presentation
import pandas as pd
import matplotlib.pyplot as plt
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

dict_days = dict(zip(['MON', 'TUE', 'WED', 'THU', 'FRI', 'SAT', 'SUN'], 
                     ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday', 'Saturday', 'Sunday']
                    ))

# Workload per day (by Week)
# Import Data

# Volumes per day
df_day = pd.read_csv('volumes per day.csv', index_col=0)
# Weeks
LIST_WEEKS = df_day['WEEK'].unique()
df_day.head()

def analysis_week(df_day, WEEK):
    
    # Filter Scope
    df_plot = df_day[df_day['WEEK'] ==WEEK].copy()

    ''' Calculate KPIs for Comments'''
    # Lines per orders
    df_plot['LINES/ORDER'] = df_plot['LINES']/df_plot['ORDERS']
    avg_ratio = '{:.2f} lines/order'.format(df_plot['LINES/ORDER'].mean())
    max_ratio = '{:.2f} lines/order'.format(df_plot['LINES/ORDER'].max())

    # Maximum Day Lines
    busy_day = dict_days[df_plot.set_index('DAY')['LINES'].idxmax()]
    max_lines = '{:,} lines'.format(df_plot['LINES'].max())

    # Total Workload
    total_lines = '{:,} lines'.format(df_plot['LINES'].sum())

    # Bar Plot: Orders/Lines
    fig, ax = plt.subplots(figsize=(12, 6))
    df_plot.plot.bar(figsize=(8, 6), edgecolor='black', x='DAY', y=['ORDERS', 'LINES'], 
                     color=['tab:blue', 'tab:orange'], legend= True, ax = ax)
    plt.xlabel('DAY', fontsize = 12)
    plt.title('Workload per day (Lines/day)', fontsize = 12)
    #  plt.show()

    # Save plot
    filename = WEEK + ".png"
    path_plot = filename
    fig.savefig(path_plot, dpi=fig.dpi)
    
    return avg_ratio, max_ratio, busy_day, max_lines, total_lines

# Number of Lines/Order (by Week)
# Importing Data
# Volumes per day
df_lior = pd.read_csv('lines per day.csv', index_col=0)
COLS_IN = list(df_lior.columns[0:8])
df_lior.reset_index(inplace = True)
df_lior.head()

def plot_split(df_lior):
    
    # Bar Plot: split per lines per orders
    fig, ax = plt.subplots(figsize=(12, 6))
    df_lior.plot.bar(figsize=(10, 6), edgecolor='black', x='WEEK', y=COLS_IN[0], color='tab:blue', legend= True, ax = ax)
    df_lior.plot.bar(figsize=(10, 6), edgecolor='black', x='WEEK', y=COLS_IN[1], color='tab:red', legend= True, ax = ax)
    df_lior.plot.bar(figsize=(10, 6), edgecolor='black', x='WEEK', y=COLS_IN[2], color='tab:orange', legend= True, ax = ax)
    df_lior.plot.bar(figsize=(10, 6), edgecolor='black', x='WEEK', y=COLS_IN[3], color='darkblue', legend= True, ax = ax)
    df_lior.plot.bar(figsize=(10, 6), edgecolor='black', x='WEEK', y=COLS_IN[4], color='brown', legend= True, ax = ax)
    df_lior.plot.bar(figsize=(10, 6), edgecolor='black', x='WEEK', y=COLS_IN[5], color='grey', legend= True, ax = ax)
    plt.xlabel('Week', fontsize = 14)
    plt.ylabel('Number of Orders', fontsize = 14)
    plt.title('Split of orders by number of lines/order', fontsize = 14)
#     plt.show()

    # Save plot
    filename = 'SPLIT' + ".png"
    path_plot = filename
    fig.savefig(path_plot, dpi=fig.dpi)

    # Analysis
    orders = df_lior.sum(axis = 1).sum()
    total_orders = '{:,} orders'.format(orders)
    
    # Split full month
    df_an = pd.DataFrame(100 * df_lior.set_index('WEEK').sum(axis = 0)/orders)
    df_an.columns = ['%']
    LIST_ANALYSIS = []
    for l in df_an.index[0:3]:
        LIST_ANALYSIS.append('{}% of orders with {} line(s) per order'.format(df_an.loc[l, '%'].round(1), l))
        
        
    return total_orders, LIST_ANALYSIS

# Creating PowerPoint
# Create the presentation object
prs = Presentation()

# Layout Choice for the introduction slide
image_slide_layout = prs.slide_layouts[5]
page = 1

'''Slide Introduction'''
# Create the slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 32, 96)
# Add Title
title.text = "WAREHOUSE WORKLOAD ANALYSIS"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
# Add Subtitle
subtitle.text = 'Orders/day for the last {} weeks'.format(len(LIST_WEEKS))
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

'''Slides Analysis per Week'''
for WEEK in LIST_WEEKS:
    
    # Create Plot
    avg_ratio, max_ratio, busy_day, max_lines, total_lines = analysis_week(df_day, WEEK)

    # Create a slide
    slide = prs.slides.add_slide(image_slide_layout)
    shapes = slide.shapes
    # Create Title
    title_shape = shapes.title
    title_shape.text = 'Warehouse Workload ({})'.format(WEEK)
    # Add Image    
    left = Inches(0.75)
    top = Inches(1.25)
    pic = slide.shapes.add_picture('{}.png'.format(WEEK), left, top, height=Inches(4.5))

    # Build the Text Box
    left = Inches(0.75)
    top = Inches(1.5) + Inches(4)
    width = Inches(9)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    # Title
    p = tf.add_paragraph()
    p.text = 'Analysis'
    p.font.size = Pt(18)
    # First bullet point
    p = tf.add_paragraph()
    p.text = '• {} have been prepared during the week'.format(total_lines)
    p.level = 1
    # Second bullet point
    p = tf.add_paragraph()
    p.text = '• {} has been the busiest day with {} prepared'.format(busy_day, max_lines)
    p.level = 1
    # Third bullet point
    p = tf.add_paragraph()
    p.text = '• {} on average with a maximum of {}'.format(avg_ratio, max_ratio)
    p.level = 1
    
    # Add Pages
    txBox = slide.shapes.add_textbox(Inches(9), Inches(6.75), Inches(1), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = str('{}/{}'.format(page, len(LIST_WEEKS) + 1 ))
    p.font.size = Pt(15)
    page += 1
    
# Add Analysis lines
total_orders, LIST_ANALYSIS = plot_split(df_lior)
# Create a slide
slide = prs.slides.add_slide(image_slide_layout)
shapes = slide.shapes
# Create Title
title_shape = shapes.title
title_shape.text = 'Order Profile'
# Add Image
left = Inches(0.75)
top = Inches(1.25)
pic = slide.shapes.add_picture('{}.png'.format('SPLIT'), left, top, height=Inches(4.5))
# Build the Text Box
left = Inches(0.75)
top = Inches(1.5) + Inches(4)
width = Inches(9)
height = Inches(2)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = '{} prepared'.format(total_orders)
p.font.size = Pt(18)
for l in LIST_ANALYSIS:
    # First bullet point
    p = tf.add_paragraph()
    p.text = '• {}'.format(l)
    p.level = 1
# Add Pages
txBox = slide.shapes.add_textbox(Inches(9), Inches(6.75), Inches(1), Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = str('{}/{}'.format(page, len(LIST_WEEKS) +1))
p.font.size = Pt(15)
page += 1

# Save
prs.save('Warehouse Workload Report.pptx')